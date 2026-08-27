"""One-off: update the deck's evidence claims and fill in the notes pane.

Writes a new file rather than editing in place, so the original deck stays
exactly as it was. Text replacements are per-run so the existing formatting
survives.
"""

from pptx import Presentation

SRC = "custody_chain_rl_stanford_v2_2.pptx"
DST = "custody_chain_rl_stanford_v2_3.pptx"

# (slide index, shape index, old text, new text)
EDITS = [
    (8, 9,
     "Unlike MIMIC, guardrail is populated: go-arounds and ATC interventions are coded refusals.",
     "Unlike MIMIC, guardrail is populated: 47% of 100 published reports carry a coded refusal."),
    (8, 15,
     "A first pass admits nothing. Review is an amendment that re-seals — visible, not silent.",
     "A first pass admits nothing — 0 of 100. Review is an amendment that re-seals, visible not silent."),
    (11, 8,
     "Research prototype. End to end on synthetic data, plus runs on MIMIC-IV and NASA ASRS. "
     "No component has been run against an operating facility's own custody ledger",
     "Research prototype. End to end on synthetic data, MIMIC-IV, and 100 published NASA ASRS "
     "reports. No component has been run against an operating facility's own custody ledger"),
]

NOTES = {
1: """Regulated facilities produce the highest-integrity training data in the building.
Collected under legal obligation, never read again. I want to argue that IS the dataset.

Speaking personally. Everything here is public or synthetic.""",

2: """The claim, stated hard: the field audits DECISIONS. Nobody asks whether the RECORD is admissible.

Update-auditing seals the training run. Statistical guarantees hold GIVEN the data.
Both assume the record is what it claims to be. I make that assumption checkable — and refuse
to train when it fails.

Three legs: the record is the dataset / integrity gates training / the reading is in the chain.""",

3: """Take the objection head-on before anyone raises it.

Three places a model sits: policy (VLA), planner (SayCan), and READER OF RECORD.
The first two are this room's work, already done. The third is where the problem is.

Reframe: post-training runs on verifiable rewards — math, code, unit tests.
Physical operations has no unit test. My claim: the custody record IS the verifier.
An episode with sealed evidence before and after is a checkable outcome.""",

4: """State the prior work first, because this room wrote most of it.

Seldonian / high-confidence SPI. OPE attribution. Proof-of-learning, in-toto, SLSA.
Offline RL in healthcare and why it goes wrong.

THE GAP: all of it assumes the dataset is admissible. Admissibility comes BEFORE validity.
The ask is genuine — if this has been formalised, I want the citation.""",

5: """One JSON object, two readers: an auditor and an RL algorithm. Same bytes.

Three things to notice.
- The seal covers the PREVIOUS seal, so an edit is locatable by index.
- Cost is not negative reward. A constraint violation cannot be traded against outcome.
- A step with no evidence FAILS validation. An action nobody can audit is not training data.
  That last rule is the line between a custody ledger and an ordinary offline RL dataset.""",

6: """Five verbs, and every one can refuse.

verify walks the chain and reports a bad seal separately from a bad link — they mean
different things to an investigator.

The exit codes are the point. A release pipeline can refuse to promote a policy that has
not earned promotion. Governance as a build step, not a review meeting.""",

7: """DEMO. Two minutes.

1. Reviewer blank -> ledger empty. The first pass admits nothing. That is the design.
2. Type a reviewer -> STILL empty. Review is necessary, not sufficient.
   Be honest here: that check used to be an 'elif', which meant passing a reviewer name
   skipped it entirely. My own test file was written to prevent exactly that — but it
   tested the helper, not the gate that used it. A keyword matcher never surfaced it.
   It took a model reading for meaning to generate enough real doubt to expose the hole.
3. Tamper one reward -> that card turns rust, verify names the index, fit refuses, exit 1.
4. Delete an episode -> bad LINK, not bad seal. Different failure, different meaning.

The refusal is the feature.""",

8: """'Why did you learn this?' answered as a list, not a story.

The estimator is deliberately boring: first-visit Monte Carlo. A reviewer can recompute
any value by hand.

Because the estimate is linear in its episodes, attribution is EXACT — no influence-function
estimation. Expect the Koh & Liang question. Their answer is approximate for a trained model;
mine is exact and free, and that is the payoff for choosing a boring estimator.""",

9: """The premise says facilities hold the data. True and misleading — they hold it as PROSE.

If a model decides what the action was, the model joins the causal chain behind the policy.
So it is sealed in: model, weights digest, prompt hash, source hash, review status.
The name is not the identity. The digest is.

REAL NUMBERS, 100 published NASA reports, two Report Sets:
- guardrail populated on 47% — the field MIMIC had to leave empty
- keyword baseline classified 4% of steps; the model classified all of them (608 -> 0)
- the model flagged 81 places it had to guess — semantic doubt, not missing regex
- admitted on first pass: 0

Note the control: guardrail stays at 47% for BOTH readers, because it comes from NASA's
coded fields, not from the model. Swap the reader and that number does not move.

Three boundaries: never extracts reward, never extracts safety claims, unreviewed is
inadmissible.""",

10: """Least novel component and I will say so. Staged autonomy is how avionics releases already work.

Four successes out of four is a raw rate of 100% and a Wilson lower bound of 51%.
The lower bound is the honest reading. Raw rates promote skills on lucky samples.

Run the gate BACKWARDS: how many more episodes at the observed rate before this skill can be
promoted? That turns the same bound into a planning instrument, which is the question an
operations director actually has.""",

11: """The sharpest limitation, found by running my own pipeline.

MIMIC-IV, 120-stay sepsis cohort: 503 learned values, 81 states, MEDIAN ONE observation
per value. That is the empirical form of Gottesman et al. 2019 — an offline estimator
returns a confident-looking number for a state-action pair clinicians almost never chose.

Two fields honestly: evidence is real, guardrail is empty. MIMIC does not log 'a protocol
blocked this action'. Inventing those events would fabricate the one thing this protects.

Worth adding: on real ASRS prose the pipeline broke three times — null timestamps,
a truncated model reply, and steps with no evidence. Every one was invisible on the
fixture, because I wrote the fixture and I wrote it to pass. The truncation one matters
most: a cut-off reading seals a record that looks perfectly valid and is missing the end
of the story.""",

12: """Say the limit yourself, before anyone says it for you.

Chain integrity is NOT causal validity. A perfectly sealed ledger can still support a
wrong policy through unmeasured confounding.

This makes the record trustworthy. It does not make the inference correct — and the
evidence bundle says so in its own text, so nobody downstream can quote it as a safety case.

Status: research prototype. Never run against an operating facility's own ledger.""",

13: """Behind an airgap this stops being a design choice.

No cloud model, no vendor fleet learning, no telemetry egress. The only training data that
will ever exist is what the facility recorded itself. Every competing approach assumes data
leaves the building.

The good argument: a one-way channel has NO ACKNOWLEDGEMENT. The receiver cannot ask
'did I get everything?' Continuity of the hash links proves completeness with no return
channel, and a broken link localises the gap.

For a utility audience: substation LAN, NERC CIP electronic security perimeter, data diodes
are already standard OT equipment. And multi-site federation without moving data.""",

14: """The direction I want to argue about.

A guardrail halt is not an outcome label. It is a COUNTERFACTUAL: the policy wanted X and
the envelope stopped it. That is a labelled unsafe-action datapoint, and an ordinary offline
RL dataset can never contain one — ordinary datasets record what happened, not what was blocked.

Learn the envelope, keep enforcement human and fixed. Flag proposals before they reach the
barrier. And measure drift toward the envelope over time as an early warning for reward
hacking — a partial answer to my own open problem, at no extra collection cost.""",

15: """Four open problems, plainly. These are why this is a proposal, not a result.

1. Reward hacking against audit metrics — if the record is the reward source, a policy can
   learn to produce well-formed records rather than good outcomes.
2. Evidence retention cost — linking updates is cheap, retaining sensor evidence is not.
3. Certification power — the gate cannot speak to conditions absent from the record.
4. Extraction ontology — accuracy is the easy half; choosing the action vocabulary is a
   value judgement made by a model, reviewed under time pressure, invisible to the chain.

Three asks: prior work, a real ledger to run the verifier over read-only, and please attack
the envelope idea.""",

16: """Thank you. Questions, disagreements, and collaborations — in that order.

If one thing lands: offline RL and update-auditing both assume the dataset is admissible.
When the dataset is prose read by a model, that assumption is false unless the reading is
sealed, reviewable, and refused until checked.""",
}


def main() -> None:
    prs = Presentation(SRC)

    for slide_i, shape_i, old, new in EDITS:
        tf = prs.slides[slide_i].shapes[shape_i].text_frame
        for para in tf.paragraphs:
            joined = "".join(r.text for r in para.runs)
            if joined.strip() != old.strip():
                continue
            para.runs[0].text = new
            for r in para.runs[1:]:
                r.text = ""
            print(f"slide {slide_i+1} shape {shape_i}: updated")
            break
        else:
            print(f"slide {slide_i+1} shape {shape_i}: NO MATCH — check by hand")

    for num, text in NOTES.items():
        prs.slides[num - 1].notes_slide.notes_text_frame.text = text
    print(f"notes written for {len(NOTES)} slide(s)")

    prs.save(DST)
    print(f"\nsaved {DST}")


if __name__ == "__main__":
    main()
